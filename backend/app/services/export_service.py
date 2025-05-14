import os
import json
import logging
import tempfile
import plotly.graph_objects as go
from typing import List, Dict, Any, Optional, Union, Tuple
import io
import base64
from pathlib import Path
import time

import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd

from app.core.config import settings
from app.schemas.export import ExportRequest, SlideRequest
from app.services.file_service import FileService

logger = logging.getLogger(__name__)

class ExportService:
    def __init__(self):
        # Create an export directory
        self.export_dir = os.path.join(settings.upload_dir, "exports")
        os.makedirs(self.export_dir, exist_ok=True)
        
        # Initialize file service
        self.file_service = FileService()
    
    async def export_chart(self, fig_json: Dict[str, Any], format: str = "png", width: Optional[int] = None, height: Optional[int] = None, scale: float = 1.0, filename: Optional[str] = None) -> str:
        """
        Export a chart as PNG or SVG.
        """
        try:
            # Create figure from JSON
            fig = go.Figure(fig_json)
            
            # Set width and height if provided
            if width and height:
                fig.update_layout(width=width, height=height)
            
            # Generate a filename if not provided
            if not filename:
                filename = f"chart_{int(time.time())}"
            
            # Ensure filename has correct extension
            if not filename.lower().endswith(f".{format}"):
                filename = f"{filename}.{format}"
            
            # Create the full output path
            output_path = os.path.join(self.export_dir, filename)
            
            # Write to file
            if format == "png":
                pio.write_image(fig, output_path, scale=scale)
            elif format == "svg":
                pio.write_image(fig, output_path)
            else:
                raise ValueError(f"Unsupported format: {format}")
            
            return output_path
        
        except Exception as e:
            logger.error(f"Error exporting chart: {str(e)}")
            raise
    
    async def export_slide(self, fig_json: Dict[str, Any], title: str, caption: Optional[str] = None, filename: Optional[str] = None) -> str:
        """
        Export a PowerPoint slide with a chart.
        """
        try:
            # Create figure from JSON
            fig = go.Figure(fig_json)
            
            # Export chart as PNG temporarily
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                pio.write_image(fig, tmp.name, scale=2.0)
                tmp_path = tmp.name
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Add a slide
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set the title
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Add the image
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            pic = slide.shapes.add_picture(tmp_path, left, top, width=width)
            
            # Add caption if provided
            if caption:
                left = Inches(1)
                top = Inches(5.5)
                width = Inches(8)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = caption
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.italic = True
            
            # Generate a filename if not provided
            if not filename:
                filename = f"presentation_{int(time.time())}.pptx"
            
            # Ensure filename has correct extension
            if not filename.lower().endswith(".pptx"):
                filename = f"{filename}.pptx"
            
            # Create the full output path
            output_path = os.path.join(self.export_dir, filename)
            
            # Save the presentation
            prs.save(output_path)
            
            # Clean up the temporary image
            os.unlink(tmp_path)
            
            return output_path
        
        except Exception as e:
            logger.error(f"Error exporting slide: {str(e)}")
            raise
    
    async def export_data(self, file_id: str, format: str = "csv") -> Tuple[str, str]:
        """
        Export a dataset as CSV or Excel.
        """
        try:
            # Find the file path from ID
            file_path = self.file_service.get_file_path_from_id(file_id)
            if not file_path:
                raise ValueError("File not found")
            
            # Generate the output filename
            filename = f"export_{os.path.basename(file_path)}"
            
            # Check if format needs to be changed
            if format == "csv" and not filename.lower().endswith(".csv"):
                filename = f"{os.path.splitext(filename)[0]}.csv"
            elif format == "excel" and not filename.lower().endswith(".xlsx"):
                filename = f"{os.path.splitext(filename)[0]}.xlsx"
            
            # Create the full output path
            output_path = os.path.join(self.export_dir, filename)
            
            # Read the dataset
            if file_path.lower().endswith('.csv'):
                df = pd.read_csv(file_path)
            elif file_path.lower().endswith('.xlsx'):
                df = pd.read_excel(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_path}")
            
            # Export the dataset
            if format == "csv":
                df.to_csv(output_path, index=False)
                media_type = "text/csv"
            elif format == "excel":
                df.to_excel(output_path, index=False)
                media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            else:
                raise ValueError(f"Unsupported export format: {format}")
            
            return output_path, media_type
        
        except Exception as e:
            logger.error(f"Error exporting data: {str(e)}")
            raise